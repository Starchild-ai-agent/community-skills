#!/usr/bin/env python3
"""
doc-to-markdown: Universal document-to-Markdown converter.
Supports PDF, Word (.docx), PowerPoint (.pptx), Excel (.xlsx), CSV, HTML, and text files.
"""

import sys
import os
import json
import argparse
import csv
import io


def convert_pdf(file_path):
    """Convert PDF to Markdown using pdfplumber."""
    import pdfplumber

    pages = []
    table_count = 0
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            parts = [f"## Page {i}"]
            text = page.extract_text()
            if text:
                parts.append(text.strip())

            tables = page.extract_tables()
            for table in tables:
                if not table or len(table) < 2:
                    continue
                table_count += 1
                md_table = _table_to_markdown(table)
                parts.append(md_table)

            pages.append("\n\n".join(parts))

    return "\n\n---\n\n".join(pages), {"pages": len(pdf.pages), "tables": table_count}


def convert_docx(file_path):
    """Convert Word document to Markdown using python-docx."""
    from docx import Document
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

    doc = Document(file_path)
    parts = []
    table_count = 0

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # Find corresponding paragraph object
            for para in doc.paragraphs:
                if para._element is element:
                    style_name = (para.style.name or "").lower() if para.style else ""
                    text = para.text.strip()
                    if not text:
                        break

                    if "heading 1" in style_name or "title" in style_name:
                        parts.append(f"# {text}")
                    elif "heading 2" in style_name:
                        parts.append(f"## {text}")
                    elif "heading 3" in style_name:
                        parts.append(f"### {text}")
                    elif "heading 4" in style_name:
                        parts.append(f"#### {text}")
                    elif "heading 5" in style_name:
                        parts.append(f"##### {text}")
                    elif "heading 6" in style_name:
                        parts.append(f"###### {text}")
                    elif "list" in style_name or style_name.startswith("list"):
                        parts.append(f"- {text}")
                    else:
                        parts.append(text)
                    break
        elif tag == "tbl":
            for table in doc.tables:
                if table._element is element:
                    table_count += 1
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        parts.append(_table_to_markdown(rows))
                    break

    return "\n\n".join(parts), {"tables": table_count}


def convert_pptx(file_path):
    """Convert PowerPoint to Markdown using python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    table_count = 0

    for i, slide in enumerate(prs.slides, 1):
        parts = [f"## Slide {i}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level if para.level else 0
                    indent = "  " * level
                    # Check if it's a title
                    if shape == slide.shapes.title:
                        parts[0] = f"## Slide {i}: {text}"
                    else:
                        parts.append(f"{indent}- {text}")

            if shape.has_table:
                table_count += 1
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    parts.append(_table_to_markdown(rows))

        slides.append("\n\n".join(parts))

    return "\n\n---\n\n".join(slides), {"slides": len(prs.slides), "tables": table_count}


def convert_xlsx(file_path):
    """Convert Excel to Markdown using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True, read_only=True)
    sheets = []
    table_count = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            str_row = [str(cell) if cell is not None else "" for cell in row]
            # Skip completely empty rows
            if any(c.strip() for c in str_row):
                rows.append(str_row)

        if rows:
            table_count += 1
            parts = [f"## Sheet: {sheet_name}", _table_to_markdown(rows)]
            sheets.append("\n\n".join(parts))

    wb.close()
    return "\n\n---\n\n".join(sheets), {"sheets": len(wb.sheetnames), "tables": table_count}


def convert_csv(file_path):
    """Convert CSV to Markdown table."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return "", {"rows": 0}

    md = _table_to_markdown(rows)
    return md, {"rows": len(rows)}


def convert_html(file_path):
    """Convert HTML to Markdown using BeautifulSoup."""
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()

    soup = BeautifulSoup(html, "html.parser")

    # Remove script and style
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    parts = []
    table_count = 0

    def process_element(elem):
        nonlocal table_count
        if isinstance(elem, str):
            text = elem.strip()
            if text:
                parts.append(text)
            return

        name = elem.name

        if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(name[1])
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"{'#' * level} {text}")
        elif name == "p":
            text = elem.get_text(strip=True)
            if text:
                parts.append(text)
        elif name == "li":
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"- {text}")
        elif name == "blockquote":
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"> {text}")
        elif name == "code":
            text = elem.get_text()
            if text:
                parts.append(f"`{text}`")
        elif name == "pre":
            text = elem.get_text()
            if text:
                parts.append(f"```\n{text}\n```")
        elif name == "table":
            rows = []
            for tr in elem.find_all("tr"):
                cells = tr.find_all(["td", "th"])
                rows.append([c.get_text(strip=True) for c in cells])
            if rows:
                table_count += 1
                parts.append(_table_to_markdown(rows))
        elif name in ("ul", "ol", "div", "section", "article", "main", "body", "header", "footer", "nav"):
            for child in elem.children:
                process_element(child)
        else:
            text = elem.get_text(strip=True)
            if text:
                parts.append(text)

    process_element(soup.body or soup)
    return "\n\n".join(parts), {"tables": table_count}


def convert_text(file_path):
    """Convert text file to Markdown (pass-through with cleanup)."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return content.strip(), {"chars": len(content)}


def _table_to_markdown(rows):
    """Convert a list of rows (list of cell strings) to a Markdown table."""
    if not rows:
        return ""

    # Normalize column count
    max_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < max_cols:
            r.append("")

    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []

    lines = []
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in body:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


CONVERTERS = {
    ".pdf": ("pdf", convert_pdf),
    ".docx": ("word", convert_docx),
    ".pptx": ("powerpoint", convert_pptx),
    ".xlsx": ("excel", convert_xlsx),
    ".csv": ("csv", convert_csv),
    ".html": ("html", convert_html),
    ".htm": ("html", convert_html),
    ".txt": ("text", convert_text),
    ".md": ("text", convert_text),
    ".rst": ("text", convert_text),
    ".log": ("text", convert_text),
}


def main():
    parser = argparse.ArgumentParser(
        description="Convert documents to Markdown"
    )
    parser.add_argument("input", help="Input file path")
    parser.add_argument("--output", "-o", help="Output file path (default: stdout)")
    parser.add_argument(
        "--format",
        choices=["markdown", "json"],
        default="markdown",
        help="Output format (default: markdown)",
    )
    args = parser.parse_args()

    file_path = args.input
    if not os.path.exists(file_path):
        error = {"success": False, "error": f"File not found: {file_path}"}
        print(json.dumps(error))
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in CONVERTERS:
        supported = ", ".join(sorted(CONVERTERS.keys()))
        error = {
            "success": False,
            "error": f"Unsupported file extension: {ext}",
            "supported": supported,
        }
        print(json.dumps(error))
        sys.exit(1)

    format_name, converter = CONVERTERS[ext]

    try:
        markdown, meta = converter(file_path)
    except Exception as e:
        error = {"success": False, "error": f"Conversion failed: {str(e)}", "format": format_name}
        print(json.dumps(error))
        sys.exit(1)

    meta["format"] = format_name
    meta["success"] = True

    if args.format == "json":
        result = {"success": True, "markdown": markdown, **meta}
        output = json.dumps(result, ensure_ascii=False, indent=2)
        if args.output:
            with open(args.output, "w", encoding="utf-8") as f:
                f.write(output)
        else:
            print(output)
    else:
        if args.output:
            with open(args.output, "w", encoding="utf-8") as f:
                f.write(markdown)
            meta_str = " ".join(f"{k}={v}" for k, v in meta.items() if k != "success")
            print(f"Converted {format_name} → {args.output} ({meta_str})", file=sys.stderr)
        else:
            print(markdown)


if __name__ == "__main__":
    main()
